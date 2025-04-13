from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

# Function to set background color
def set_slide_background(slide, color=RGBColor(255, 255, 255)):  # Default white background
    slide_background = slide.background
    fill = slide_background.fill
    fill.solid()
    fill.fore_color.rgb = color

# Function to add a slide with title, content, and apply global styles
def add_styled_slide(prs, title, content, img_path=None):
    slide_layout = prs.slide_layouts[1]  # Title and Content layout
    slide = prs.slides.add_slide(slide_layout)

    # Set background color (light gray here)
    set_slide_background(slide, RGBColor(242, 242, 242))

    # Title and body content
    title_placeholder = slide.shapes.title
    body_shape = slide.shapes.placeholders[1]
    title_placeholder.text = title

    # Apply title style (font size, color, bold, etc.)
    set_font_style(slide, is_title=True)

    # Add content text
    tf = body_shape.text_frame
    tf.text = content
    
    # Apply body text style
    set_font_style(slide, is_title=False)

    # Add image if image path is provided
    if img_path and os.path.exists(img_path):
        left = Inches(1)
        top = Inches(2.5)
        pic = slide.shapes.add_picture(img_path, left, top, height=Inches(3))

# Function to set font style for title and body text
def set_font_style(slide, is_title=True):
    if is_title:
        title_shape = slide.shapes.title
        title_text_frame = title_shape.text_frame
        for paragraph in title_text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.size = Pt(36)  # Font size for title
                run.font.bold = True    # Bold font for title
                run.font.color.rgb = RGBColor(0, 0, 0)  # Black color for title text
                paragraph.alignment = PP_ALIGN.CENTER
    else:
        body_shape = slide.shapes.placeholders[1]
        body_text_frame = body_shape.text_frame
        for paragraph in body_text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.size = Pt(18)  # Font size for body text
                run.font.color.rgb = RGBColor(0, 0, 0)  # Black color for body text
            paragraph.alignment = PP_ALIGN.LEFT  # Left alignment for content

# Create a presentation object
prs = Presentation()

# Slide 1: Title Slide
add_styled_slide(prs, 
                 "Predicting Winners in Chatbot Arena", 
                 "How do humans do it? \n Scott Gross, 4/13/2025")

# Slide 2: Project Overview
content = """Goal: Build a predictive classifier for human preference in chatbot responses.
Model Predicts: Which of two model-generated replies is preferred based on latent criteria (subtle, often unconscious rules) that humans actually use to judge quality.
Question: Are syntactic features (length, punctuation, word count) or semantic understanding (embeddings, originality) more predictive? Is something more required?"""
add_styled_slide(prs, "Project Overview", content)

# Slide 3: Data Exploration & Preprocessing
content = """Data: Excel spreadsheet → CSV 
Prompt | response_A | response_B | LABEL (0 = A wins, 1 = B wins, 2 = Tie).
Feature Engineering (additional columns):
Length, Word Count, Verbosity, Punctuation.
Class Balance: No significant class imbalance; slight underrepresentation of ties."""
add_styled_slide(prs, "Data Exploration & Preprocessing", content)

# Slide 4: Syntactic Approaches
content = """Initial Hypothesis: Simple surface features like length and punctuation should explain preference.
Tools Used: 
- Punctuation count, verbosity score, word count, length difference.
Model: Random Forest.
Results: Weak signal (~39% accuracy), highlighting limitations of surface-level features."""
add_styled_slide(prs, "Syntactic Approaches", content)

# Slide 5: Stochastic Approaches
content = """Embedding Approach: Using cosine similarity and sentence embeddings.
Tools: Sentence-transformers, Logistic Regression / MLP / XGBoost"""
add_styled_slide(prs, "Stochastic Approaches", content, img_path='cosine_results_1.jpg')

# Slide 6: Results - Initial Approaches
content = """Accuracy from Random Forest: ~39%
Accuracy from XGBoost: ~38%
Accuracy from Logistic Regression: ~57.5%
Key Insight: Surface-level features and semantic similarity alone aren't enough to predict human preferences."""
add_styled_slide(prs, "Results - Initial Approaches", content)

# Slide 7: Realization - Syntactic == Stochastic
content = """Insight: The embeddings (semantic representations) were capturing patterns that surface-level features already partially represented.
Conclusion: Stochastic approaches were too similar to syntactic features in this context."""
add_styled_slide(prs, "Realization - Syntactic == Stochastic", content)

# Slide 8: Introducing Heuristics
content = """Heuristics Added:
- Keyword presence and salience.
- Factuality and depth of responses.
- Originality scoring."""
add_styled_slide(prs, "Introducing Heuristics", content, img_path='originality_results.jpg')

# Slide 9: Results from Heuristic Tests
content = """Logistic Regression: 57.5% accuracy.
Random Forest: 36.0% accuracy.
XGBoost: 38.0% accuracy.
Conclusion: Logistic Regression performed best; however, further improvements are needed."""
add_styled_slide(prs, "Results from Heuristic Tests", content)

# Slide 10: The MVAT Approach
content = """MVAT (Multi-Vector Assessment Tree): A layered approach to scoring.
- Integrates multiple linguistic features into a decision tree for more nuanced prediction."""
add_styled_slide(prs, "The MVAT Approach", content)

# Slide 11: Results from ML Algorithms
content = """Results from 3 different ML algorithms:
- Logistic Regression
- Random Forest
- XGBoost"""
add_styled_slide(prs, "Results from ML Algorithms", content, img_path='syntax_results.jpg')

# Slide 12: Let's Train an LLM Locally
content = """The decision to train a local model due to the semantic understanding required to better capture human preference. Training an LLM locally allows for more direct optimization for the task at hand."""
add_styled_slide(prs, "Let's Train an LLM Locally", content)

# Save the PowerPoint presentation
prs.save("chatbot_arena_presentation_styled.pptx")
print("Styled PowerPoint presentation generated successfully!")



